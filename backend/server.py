from fastapi import FastAPI, HTTPException, Depends, File, UploadFile, Form, Body
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from pydantic_core import core_schema
from typing import List, Optional, Dict, Any
import os
import json
import uuid
from datetime import datetime
import google.generativeai as genai
import PyPDF2
from docx import Document
from pptx import Presentation
import io
from supabase import create_client, Client
from pymongo import MongoClient
from motor.motor_asyncio import AsyncIOMotorClient
from bson.objectid import ObjectId
import asyncio
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Initialize FastAPI app
app = FastAPI(title="PaperGenius API", version="1.0.0")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"], # Explicitly allow frontend origin
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Environment variables
MONGO_URL = os.environ.get('MONGO_URL', 'mongodb://localhost:27017')
DB_NAME = os.environ.get('DB_NAME', 'papergenius')
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
SUPABASE_URL = os.environ.get('SUPABASE_URL')
SUPABASE_ANON_KEY = os.environ.get('SUPABASE_ANON_KEY')

# Initialize services
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

# Initialize Supabase (optional for now as we're using MongoDB)
supabase = None
if SUPABASE_URL and SUPABASE_ANON_KEY:
    try:
        supabase = create_client(SUPABASE_URL, SUPABASE_ANON_KEY)
        print("Supabase client initialized successfully")
    except Exception as e:
        print(f"Warning: Could not initialize Supabase client: {e}")
else:
    print("Warning: Supabase credentials not found in environment")

# MongoDB connection
client = AsyncIOMotorClient(MONGO_URL)
db = client[DB_NAME]

# Pydantic models
class FolderCreate(BaseModel):
    name: str
    type: str  # "unit-wise" or "syllabus"
    user_id: str

class UnitQuestionSpec(BaseModel):
    unit_name: str 
    num_questions: int
    # Optional: Add other per-unit specifics like difficulty, question_type if needed later
    # E.g., question_type: Optional[str] = "MCQ" 
    # E.g., difficulty: Optional[str] = "Medium"

class SectionDetail(BaseModel):
    section_name: str
    section_type: str # E.g., "Answer All Questions", "Answer N from Section"
    questions_type: str # E.g., "MCQ", "Short Answer", "Long Answer", "Fill in the Blanks", "True/False"
    total_questions: int
    questions_to_be_answered: int
    marks_for_each_question: int
    custom_instruction: Optional[str] = ""
    question_specs: List[UnitQuestionSpec]

class TemplateCreate(BaseModel):
    name: str
    description: Optional[str] = ""
    instituteType: str
    instituteName: Optional[str] = ""
    evaluation: str
    duration: int
    paper_code: Optional[str] = "" # Made optional
    total_marks: Optional[int] = 0 # Made optional, should be calculated
    sections: List[SectionDetail]
    user_id: str

class QuestionPaperGenerate(BaseModel):
    folder_id: str
    template_id: str
    selected_units: List[str]
    paper_name: str
    user_id: str

# Helper for MongoDB ObjectId
class PyObjectId(ObjectId):
    @classmethod
    def __get_validators__(cls):
        yield cls.validate

    @classmethod
    def validate(cls, v: Any) -> ObjectId:
        if isinstance(v, ObjectId):
            return v
        if ObjectId.is_valid(v):
            return ObjectId(v)
        raise ValueError("Invalid ObjectId")

    @classmethod
    def __get_pydantic_core_schema__(cls, source_type: Any, handler: Any) -> core_schema.CoreSchema:
        def validate_from_json_or_python(value: Any) -> ObjectId:
            if isinstance(value, ObjectId):
                return value
            if isinstance(value, str) and ObjectId.is_valid(value):
                return ObjectId(value)
            raise TypeError("ObjectId_expected")

        from_json_or_python_schema = core_schema.chain_schema([
            core_schema.union_schema([
                core_schema.is_instance_schema(ObjectId),
                core_schema.str_schema(),
            ]),
            core_schema.no_info_plain_validator_function(validate_from_json_or_python),
        ])

        return core_schema.json_or_python_schema(
            json_schema=core_schema.str_schema(),
            python_schema=from_json_or_python_schema,
            serialization=core_schema.plain_serializer_function_ser_schema(lambda x: str(x))
        )

# Models for displaying a single question paper
class QuestionDisplay(BaseModel):
    question: str
    type: str
    difficulty: str
    marks: int
    options: Optional[List[str]] = None
    answer: Optional[str] = None
    unit_name_source: str
    section_name_source: str

class PaperSectionDisplay(BaseModel):
    section_name: str
    section_type: str
    questions_type: str
    total_questions_in_section: int
    questions_to_be_answered_in_section: int
    marks_for_each_question_in_section: int
    custom_instruction_for_section: Optional[str] = ""
    questions: List[QuestionDisplay]
    unit_distribution_specs: List[UnitQuestionSpec] # Reusing existing model

class TemplateDetailsDisplay(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    instituteType: Optional[str] = None
    instituteName: Optional[str] = None
    evaluation: Optional[str] = None
    duration: Optional[int] = None
    paper_code: Optional[str] = None
    overall_total_marks_from_template: Optional[int] = None

class QuestionPaperDisplay(BaseModel):
    id: str  # String UUID
    mongo_id: PyObjectId = Field(alias="_id") # MongoDB ObjectId
    name: str
    user_id: str
    folder_id: str
    template_id: str
    generated_at: str
    selected_units_for_generation: List[str]
    template_details: TemplateDetailsDisplay
    paper_sections: List[PaperSectionDisplay]

    class Config:
        populate_by_name = True
        json_encoders = {ObjectId: str}
        arbitrary_types_allowed = True # For PyObjectId

# Text extraction functions
def extract_text_from_pdf(file_content: bytes) -> str:
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_docx(file_content: bytes) -> str:
    try:
        doc = Document(io.BytesIO(file_content))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"

def extract_text_from_pptx(file_content: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def extract_text_from_file(filename: str, file_content: bytes) -> str:
    """Extract text from uploaded files"""
    if filename.lower().endswith('.pdf'):
        return extract_text_from_pdf(file_content)
    elif filename.lower().endswith('.docx'):
        return extract_text_from_docx(file_content)
    elif filename.lower().endswith('.pptx'):
        return extract_text_from_pptx(file_content)
    else:
        return "Unsupported file format"

async def generate_questions_with_gemini(
    content: str, 
    num_questions: int,
    question_type_for_section: str, # E.g., "MCQ", "Short Answer"
    marks_for_each_question_in_section: int,
    # Optional: custom_instruction_for_section: Optional[str] = None,
    # Optional: section_name: Optional[str] = None
) -> List[Dict]:
    """Generate questions using Gemini AI based on specified type and marks."""
    print(f"generate_questions_with_gemini called with: num_questions={num_questions}, type='{question_type_for_section}', marks={marks_for_each_question_in_section}, content_length={len(content)}") # DEBUG
    if num_questions == 0: # DEBUG
        print("generate_questions_with_gemini: num_questions is 0, returning empty list immediately.") # DEBUG
        return [] # DEBUG
    try:
        # Verify API key is configured
        if not GEMINI_API_KEY:
            print("Warning: Gemini API key not configured. Using fallback.") # DEBUG
            return generate_fallback_questions(content, num_questions)
        
        # Test API key by making a simple request first
        model = genai.GenerativeModel('gemini-1.5-flash')
        
        # Dynamically build the JSON structure example for the prompt
        json_example_question = {
            "question": "Question text here",
            "type": question_type_for_section,
            "difficulty": "Medium", # Can be made dynamic later if needed
            "marks": marks_for_each_question_in_section
        }

        if question_type_for_section == "MCQ":
            json_example_question["options"] = ["A. Option 1", "B. Option 2", "C. Option 3", "D. Option 4"]
            json_example_question["answer"] = "A. Option 1"
        elif question_type_for_section == "True/False":
            json_example_question["options"] = ["True", "False"]
            json_example_question["answer"] = "True"
        elif question_type_for_section == "Fill in the Blanks":
            json_example_question["answer"] = "The correct answer phrase"
            # No options for fill in the blanks typically
        else: # Short Answer, Long Answer
            json_example_question["answer"] = "A model answer or key points for evaluation."
            # No options for these types

        # Prepare a concise prompt to minimize token usage
        prompt = f"""
        Based on the following educational content, generate exactly {num_questions} questions of type '{question_type_for_section}'.
        Ensure the output is a valid JSON array.
        
        Content: {content[:1500]}...
        
        Generate questions with this exact JSON structure for each question in the array:
        {json.dumps([json_example_question], indent=4)}
        
        Make questions educational and relevant to the content. For 'MCQ' or 'True/False', provide options and a single correct answer. For other types, provide a model answer.
        If the question type is 'Fill in the Blanks', the question text should clearly indicate where the blank is (e.g., using '___').
        """
        
        print(f"Gemini Prompt for {num_questions} questions of type {question_type_for_section}:\n{prompt[:500]}...\n") # DEBUG (log first 500 chars)
        response = model.generate_content(prompt)
        
        if response and response.text:
            # Try to extract JSON from the response
            response_text = response.text.strip()
            
            # Find JSON array in the response
            start_idx = response_text.find('[')
            end_idx = response_text.rfind(']') + 1
            
            if start_idx != -1 and end_idx > start_idx:
                json_str = response_text[start_idx:end_idx]
                questions = json.loads(json_str)
                
                # Validate the structure
                if isinstance(questions, list) and len(questions) > 0:
                    print(f"✅ Successfully generated {len(questions)} questions using Gemini AI")
                    questions_list = questions
            else:
                # If JSON parsing fails or returns empty list, create structured fallback
                print(f"Warning: Could not parse Gemini response as JSON or received empty list. Response text: {response_text[:500]}... Using structured fallback.") # DEBUG
                questions_list = generate_fallback_questions(content, num_questions)
            # The return is now handled after the try-except block
            
    except Exception as e:
        print(f"Gemini API Error: {str(e)}. Using fallback.") # DEBUG
        questions_list = generate_fallback_questions(content, num_questions)
    
    # Add delay here before returning, regardless of path taken
    print(f"    Pausing for 4 seconds to respect API rate limits before returning from generate_questions_with_gemini...") # DEBUG
    await asyncio.sleep(4) 
    return questions_list

def generate_fallback_questions(content: str, num_questions: int = 5) -> List[Dict]:
    """Generate fallback questions when AI is not available"""
    # Extract key topics from content
    words = content.lower().split()
    common_terms = ['algorithm', 'data', 'computer', 'programming', 'software', 'hardware', 'network', 'database', 'system', 'technology']
    found_terms = [term for term in common_terms if term in words]
    
    questions = []
    question_types = ['MCQ', 'Short', 'Long']
    difficulties = ['Easy', 'Medium', 'Hard']
    
    for i in range(min(num_questions, 10)):  # Limit to reasonable number
        q_type = question_types[i % len(question_types)]
        difficulty = difficulties[i % len(difficulties)]
        topic = found_terms[i % len(found_terms)] if found_terms else f"topic {i+1}"
        
        if q_type == 'MCQ':
            question = {
                "question": f"What is the primary concept related to {topic} in the given content?",
                "type": "MCQ",
                "difficulty": difficulty,
                "marks": 2,
                "options": [
                    f"A. Basic understanding of {topic}",
                    f"B. Advanced application of {topic}",
                    f"C. Theoretical foundation of {topic}",
                    f"D. Practical implementation of {topic}"
                ],
                "answer": f"A. Basic understanding of {topic}"
            }
        elif q_type == 'Short':
            question = {
                "question": f"Briefly explain the significance of {topic} as mentioned in the content.",
                "type": "Short",
                "difficulty": difficulty,
                "marks": 5,
                "answer": f"The content discusses {topic} as an important concept with practical applications and theoretical foundations."
            }
        else:  # Long
            question = {
                "question": f"Discuss in detail the role and importance of {topic} based on the provided content.",
                "type": "Long",
                "difficulty": difficulty,
                "marks": 10,
                "answer": f"A comprehensive discussion of {topic} should include its definition, applications, importance, and relationship to other concepts as outlined in the content."
            }
        
        questions.append(question)
    
    return questions

# API Routes
@app.get("/api/test-gemini")
async def test_gemini_api():
    """Test Gemini API connection"""
    try:
        if not GEMINI_API_KEY:
            return {"status": "error", "message": "Gemini API key not configured"}
        
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content("Hello, this is a test. Please respond with 'Hello from Gemini!'")
        
        return {
            "status": "success", 
            "message": "Gemini API is working",
            "response": response.text if response else "No response"
        }
    except Exception as e:
        return {"status": "error", "message": f"Gemini API error: {str(e)}"}

@app.get("/")
async def root():
    return {"message": "PaperGenius API is running!", "version": "1.0.0"}

@app.post("/api/folders")
async def create_folder(folder: FolderCreate):
    """Create a new folder for organizing files"""
    folder_data = {
        "id": str(uuid.uuid4()),
        "name": folder.name,
        "type": folder.type,
        "user_id": folder.user_id,
        "created_at": datetime.utcnow().isoformat(),
        "files": []
    }
    
    result = await db.folders.insert_one(folder_data)
    return {"id": folder_data["id"], "message": "Folder created successfully"}

@app.get("/api/folders/{user_id}")
async def get_user_folders(user_id: str):
    """Get all folders for a user"""
    folders = []
    async for folder in db.folders.find({"user_id": user_id}):
        folder["_id"] = str(folder["_id"])
        folders.append(folder)
    return folders

@app.post("/api/upload/{folder_id}")
async def upload_file(
    folder_id: str,
    file: UploadFile = File(...),
    unit_name: Optional[str] = Form(None)
):
    """Upload and process files"""
    try:
        # Read file content
        content = await file.read()
        
        # Extract text from file
        extracted_text = extract_text_from_file(file.filename, content)
        
        # Store file info in database
        file_data = {
            "id": str(uuid.uuid4()),
            "filename": file.filename,
            "unit_name": unit_name or "general",
            "content_type": file.content_type,
            "extracted_text": extracted_text,
            "uploaded_at": datetime.utcnow().isoformat()
        }
        
        # Update folder with new file
        await db.folders.update_one(
            {"id": folder_id},
            {"$push": {"files": file_data}}
        )
        
        return {"message": "File uploaded and processed successfully", "file_id": file_data["id"]}
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.post("/api/templates")
async def create_template(template: TemplateCreate):
    """Create a question paper template"""
    # Convert Pydantic model to dictionary to ensure all fields, including new optional ones and nested section details, are captured.
    template_dict = template.model_dump()

    template_data = {
        "id": str(uuid.uuid4()),
        "name": template_dict.get("name"),
        "description": template_dict.get("description"),
        "instituteType": template_dict.get("instituteType"),
        "instituteName": template_dict.get("instituteName"),
        "evaluation": template_dict.get("evaluation"),
        "duration": template_dict.get("duration"),
        "paper_code": template_dict.get("paper_code"), 
        "total_marks": template_dict.get("total_marks"), 
        "sections": template_dict.get("sections"), # Contains all updated section fields
        "user_id": template_dict.get("user_id"),
        "created_at": datetime.utcnow().isoformat()
    }
    
    result = await db.templates.insert_one(template_data)
    return {"id": template_data["id"], "message": "Template created successfully"}

@app.get("/api/templates/{user_id}")
async def get_user_templates(user_id: str):
    """Get all templates for a user"""
    templates = []
    async for template in db.templates.find({"user_id": user_id}):
        template["_id"] = str(template["_id"])
        templates.append(template)
    return templates

@app.post("/api/generate-paper")
async def generate_question_paper(request: QuestionPaperGenerate):
    print(f"DEBUG: generate_question_paper called with request.paper_name: '{request.paper_name}', user_id: '{request.user_id}'") # DEBUG
    """Generate question paper using AI"""
    try:
        # Get folder and template data
        folder = await db.folders.find_one({"id": request.folder_id})
        template = await db.templates.find_one({"id": request.template_id})
        
        if not folder or not template:
            raise HTTPException(status_code=404, detail="Folder or template not found")
        
        # --- Start of New Logic for Per-Unit Question Generation ---
        paper_sections_with_questions = [] # New list to store sections with their details and questions
        
        try:
            # Create a dictionary to quickly access content for selected units
            # request.selected_units still defines the overall pool of units whose content is available.
            available_unit_contents = {}
            for file_data in folder.get("files", []):
                unit_name = file_data.get("unit_name")
                if unit_name in request.selected_units:
                    if unit_name not in available_unit_contents:
                        available_unit_contents[unit_name] = ""
                    available_unit_contents[unit_name] += file_data.get("extracted_text", "") + "\n\n" # Ensure proper newlines for joining

            # Debugging prints after populating available_unit_contents
            print(f"DEBUG: request.selected_units: {request.selected_units}")
            print(f"DEBUG: available_unit_contents loaded for units: {list(available_unit_contents.keys())}")
            for unit_name_key, content_val in available_unit_contents.items():
                print(f"DEBUG: Content length for '{unit_name_key}': {len(content_val.strip()) if content_val else 0}")

            # Check if any meaningful content was actually loaded for the selected units
            meaningful_content_found = False
            if available_unit_contents:
                for content_val in available_unit_contents.values():
                    if content_val and content_val.strip(): # Check if content is not None and not just whitespace
                        meaningful_content_found = True
                        break
            
            if not meaningful_content_found:
                print(f"DEBUG: No meaningful content found. Triggering error.")
                raise HTTPException(status_code=400, detail="No meaningful content found for any of the selected units. Please check unit names and ensure files have been uploaded and processed.")

            # Process each section defined in the template
            # The 'template' variable here is a dictionary loaded from MongoDB.
            for section_detail_dict in template.get("sections", []):
                # Manually parse section_detail_dict into SectionDetail Pydantic model for validation and easy access
                section_obj = SectionDetail(**section_detail_dict)
                print(f"Processing section: {section_obj.section_name}, Type: {section_obj.section_type}, Q_Type: {section_obj.questions_type}") # DEBUG
                
                current_section_generated_questions = [] # Questions generated for the current section
                if not section_obj.question_specs: # DEBUG
                    print(f"WARNING: Section '{section_obj.section_name}' has no question_specs.") # DEBUG

                for unit_spec in section_obj.question_specs:
                    unit_name_from_template = unit_spec.unit_name
                    num_questions_for_unit = unit_spec.num_questions
                    
                    print(f"  Processing unit_spec: {unit_name_from_template}, num_questions_for_unit: {num_questions_for_unit}") # DEBUG

                    if num_questions_for_unit == 0:
                        print(f"  Skipping Gemini call for {unit_name_from_template} as num_questions_for_unit is 0.") # DEBUG
                        continue # Skip if no questions are needed for this unit

                    unit_content = available_unit_contents.get(unit_name_from_template, "")
                    if not unit_content:
                        print(f"  WARNING: No content found for unit '{unit_name_from_template}' in section '{section_obj.section_name}'. Skipping question generation for this unit.") # DEBUG
                        continue
                        print(f"Warning: No content found for unit '{unit_name_from_template}' in folder '{request.folder_id}' (specified in template), skipping question generation for this unit.")
                        continue
                    
                    print(f"    Calling generate_questions_with_gemini for unit '{unit_name_from_template}' with {num_questions_for_unit} questions.") # DEBUG
                    unit_specific_questions = await generate_questions_with_gemini(
                        content=unit_content,
                        num_questions=num_questions_for_unit,
                        question_type_for_section=section_obj.questions_type,
                        marks_for_each_question_in_section=section_obj.marks_for_each_question
                        # custom_instruction_for_section=section_obj.custom_instruction, # If generate_questions_with_gemini is updated
                        # section_name=section_obj.section_name # If generate_questions_with_gemini is updated
                    )
                    print(f"    Received {len(unit_specific_questions)} questions from generate_questions_with_gemini for unit '{unit_name_from_template}'.") # DEBUG
                    
                    for q in unit_specific_questions:
                        q["unit_name_source"] = unit_name_from_template
                        q["section_name_source"] = section_obj.section_name 
                    
                    current_section_generated_questions.extend(unit_specific_questions)
                
                # Add this section's details and its generated questions to the main list
                paper_sections_with_questions.append({
                    "section_name": section_obj.section_name,
                    "section_type": section_obj.section_type,
                    "questions_type": section_obj.questions_type,
                    "total_questions_in_section": section_obj.total_questions,
                    "questions_to_be_answered_in_section": section_obj.questions_to_be_answered,
                    "marks_for_each_question_in_section": section_obj.marks_for_each_question,
                    "custom_instruction_for_section": section_obj.custom_instruction,
                    "questions": current_section_generated_questions,
                    "unit_distribution_specs": [spec.model_dump() for spec in section_obj.question_specs] # Store specs used
                })

            print(f"Final paper_sections_with_questions: {paper_sections_with_questions}") # DEBUG
            if not paper_sections_with_questions or not any(sec.get("questions") for sec in paper_sections_with_questions):
                print("ERROR: Condition for 'No questions generated' met.") # DEBUG
                raise HTTPException(status_code=400, detail="No questions were generated for any section. Check content and template configuration.")
                
        except Exception as e:
            print(f"Error during template processing or question generation: {str(e)}")
            # Consider logging the full traceback here for debugging
            # import traceback; traceback.print_exc()
            raise HTTPException(status_code=500, detail=f"Internal error processing template or generating questions: {str(e)}")

        # --- End of New Logic ---
        
        # Create question paper
        print(f"DEBUG: About to create paper_data. request.paper_name is: '{request.paper_name}'") # DEBUG
        paper_data = {
            "id": str(uuid.uuid4()),
            "name": request.paper_name,
            "user_id": request.user_id,
            "folder_id": request.folder_id,
            "template_id": request.template_id,
            "generated_at": datetime.utcnow().isoformat(),
            "selected_units_for_generation": request.selected_units,
            
            "template_details": { # Store a snapshot of the relevant template fields from MongoDB dict
                "name": template.get("name"),
                "description": template.get("description"),
                "instituteType": template.get("instituteType"),
                "instituteName": template.get("instituteName"),
                "evaluation": template.get("evaluation"),
                "duration": template.get("duration"),
                "paper_code": template.get("paper_code"), 
                "overall_total_marks_from_template": template.get("total_marks"), # Original total marks from template if any
            },
            
            "paper_sections": paper_sections_with_questions # Structured list of sections with their questions
        }
        
        print(f"DEBUG: Attempting to save paper_data: {{'id': paper_data.get('id'), 'name': paper_data.get('name'), 'user_id': paper_data.get('user_id'), 'folder_id': paper_data.get('folder_id'), 'template_id': paper_data.get('template_id'), 'num_sections': len(paper_data.get('paper_sections', [])), 'num_total_questions': sum(len(s.get('questions', [])) for s in paper_data.get('paper_sections', []))}}") # Log key fields
        result = await db.question_papers.insert_one(paper_data)
        paper_data["_id"] = str(result.inserted_id)
        
        return paper_data
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating paper: {str(e)}")


@app.get("/api/paper/{paper_id}", response_model=QuestionPaperDisplay)
async def get_question_paper_by_id(paper_id: str):
    """Fetch a single question paper by its string ID."""
    # The 'id' field in MongoDB is the string UUID we generated.
    paper = await db.question_papers.find_one({"id": paper_id})
    if paper:
        # Pydantic model with alias for _id and PyObjectId helper should handle conversion
        return paper
    raise HTTPException(status_code=404, detail=f"Question paper with id {paper_id} not found")

@app.get("/api/papers/{user_id}")
async def get_user_papers(user_id: str):
    """Get all generated papers for a user"""
    papers = []
    async for paper in db.question_papers.find({"user_id": user_id}):
        paper["_id"] = str(paper["_id"])
        papers.append(paper)
    return papers

@app.get("/api/paper/{paper_id}")
async def get_paper_details(paper_id: str):
    """Get detailed view of a specific paper"""
    paper = await db.question_papers.find_one({"id": paper_id})
    if not paper:
        raise HTTPException(status_code=404, detail="Paper not found")
    
    paper["_id"] = str(paper["_id"])
    return paper

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
